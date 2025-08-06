import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from datetime import datetime

def criar_relatorio_ppt(all_indicators, all_dre_data, df_culturas_for_excel, nomes_cenarios, anos, all_indicators_cultura_cenarios=None):
    """
    Gera um relatório em PowerPoint com indicadores financeiros e análise por cultura.
    
    Args:
        all_indicators: Indicadores por cenário
        all_dre_data: Dados do DRE por cenário
        df_culturas_for_excel: DataFrame com dados das culturas
        nomes_cenarios: Lista com nomes dos cenários
        anos: Lista dos anos
        all_indicators_cultura_cenarios: Indicadores por cultura e cenário
    """
    try:
        # Importações necessárias
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import matplotlib.pyplot as plt
        
        # Criar nova apresentação
        prs = Presentation()
        
        # Função auxiliar para criar tabela
        def criar_tabela_slide(title_text, data_dict, anos, slide_layout_idx=1):
            slide_layout = prs.slide_layouts[slide_layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = title_text
            
            # Remover placeholder de conteúdo se existir
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                sp = content_placeholder.element
                sp.getparent().remove(sp)
            
            # Criar tabela
            rows = len(data_dict) + 1  # +1 para cabeçalho
            cols = len(anos) + 1  # +1 para nome da linha
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Cabeçalho
            table.cell(0, 0).text = "Item"
            for i, ano in enumerate(anos):
                table.cell(0, i + 1).text = str(ano)
            
            # Dados
            for i, (item_name, values) in enumerate(data_dict.items()):
                table.cell(i + 1, 0).text = item_name
                for j, value in enumerate(values):
                    if isinstance(value, (int, float)):
                        table.cell(i + 1, j + 1).text = f"R$ {value:,.0f}"
                    else:
                        table.cell(i + 1, j + 1).text = str(value)
            
            # Formatação da tabela
            for row in table.rows:
                for cell in row.cells:
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    
            # Destacar cabeçalho
            for j in range(cols):
                cell = table.cell(0, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            
            return slide

        # Função auxiliar para criar tabela de indicadores
        def criar_tabela_indicadores(title_text, indicators, anos):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = title_text
            
            # Remover placeholder de conteúdo
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                sp = content_placeholder.element
                sp.getparent().remove(sp)
            
            # Preparar dados dos indicadores
            indicadores_para_tabela = {}
            
            for key, values in indicators.items():
                if key not in ["CAGR Receita (%)", "CAGR Lucro Líquido (%)"]:
                    if isinstance(values, list):
                        if "(%)" in key or "Margem" in key:
                            indicadores_para_tabela[key] = [f"{v:.2f}%" for v in values]
                        elif "R$" in key or "Hectare" in key or "Produtividade" in key:
                            indicadores_para_tabela[key] = [f"R$ {v:,.0f}" for v in values]
                        else:
                            indicadores_para_tabela[key] = [f"{v:.2f}" for v in values]
            
            # Adicionar CAGR no final
            indicadores_para_tabela["CAGR Receita (%)"] = [f"{indicators['CAGR Receita (%)']:.2f}%"] + [""] * (len(anos) - 1)
            indicadores_para_tabela["CAGR Lucro Líquido (%)"] = [f"{indicators['CAGR Lucro Líquido (%)']:.2f}%"] + [""] * (len(anos) - 1)
            
            # Criar tabela
            rows = len(indicadores_para_tabela) + 1
            cols = len(anos) + 1
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Cabeçalho
            table.cell(0, 0).text = "Indicador"
            for i, ano in enumerate(anos):
                table.cell(0, i + 1).text = str(ano)
            
            # Dados
            for i, (indicator_name, values) in enumerate(indicadores_para_tabela.items()):
                table.cell(i + 1, 0).text = indicator_name
                for j, value in enumerate(values):
                    table.cell(i + 1, j + 1).text = str(value)
            
            # Formatação
            for row in table.rows:
                for cell in row.cells:
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Destacar cabeçalho
            for j in range(cols):
                cell = table.cell(0, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            
            return slide

        # SLIDE 1: ABERTURA
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Análise Financeira Completa do Agronegócio"
        subtitle.text = f"Indicadores por Cenário e Cultura | {anos[0]} - {anos[-1]}\nRelatório Gerado em {datetime.now().strftime('%d/%m/%Y')}"
        
        # SLIDE 2: RECEITA POR CULTURA (TABELA)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "🌾 Receita por Cultura (Ano Base)"
        
        # Remover placeholder de conteúdo
        content_placeholder = slide.placeholders[1]
        sp = content_placeholder.element
        sp.getparent().remove(sp)
        
        # Criar tabela de receitas por cultura
        rows = len(df_culturas_for_excel) + 2  # +1 cabeçalho +1 total
        cols = 4  # Cultura, Área, Receita Total, Receita/Ha
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Cabeçalho
        headers = ["Cultura", "Área (ha)", "Receita Total", "Receita/Ha"]
        for i, header in enumerate(headers):
            table.cell(0, i).text = header
        
        # Dados por cultura
        for i, (_, row) in enumerate(df_culturas_for_excel.iterrows()):
            receita_por_ha = row['Receita Total'] / row['Área (ha)'] if row['Área (ha)'] > 0 else 0
            table.cell(i + 1, 0).text = row['Cultura']
            table.cell(i + 1, 1).text = f"{row['Área (ha)']:.1f}"
            table.cell(i + 1, 2).text = f"R$ {row['Receita Total']:,.0f}"
            table.cell(i + 1, 3).text = f"R$ {receita_por_ha:,.0f}"
        
        # Linha de total
        total_row = len(df_culturas_for_excel) + 1
        table.cell(total_row, 0).text = "TOTAL"
        table.cell(total_row, 1).text = f"{df_culturas_for_excel['Área (ha)'].sum():.1f}"
        table.cell(total_row, 2).text = f"R$ {df_culturas_for_excel['Receita Total'].sum():,.0f}"
        table.cell(total_row, 3).text = f"R$ {df_culturas_for_excel['Receita Total'].sum() / df_culturas_for_excel['Área (ha)'].sum():,.0f}"
        
        # Formatação
        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(11)
        
        # Destacar cabeçalho
        for j in range(4):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Destacar total
        for j in range(4):
            cell = table.cell(total_row, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(217, 217, 217)
            cell.text_frame.paragraphs[0].font.bold = True

        # SLIDE 3: DADOS BASE (mantém texto)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "📋 Dados Base do Sistema"
        
        dados_base_text = "CONFIGURAÇÕES GERAIS:\n\n"
        dados_base_text += f"• Período de Análise: {anos[0]} a {anos[-1]} ({len(anos)} anos)\n"
        dados_base_text += f"• Número de Culturas: {len(df_culturas_for_excel)}\n"
        dados_base_text += f"• Área Total: {df_culturas_for_excel['Área (ha)'].sum():.1f} hectares\n"
        dados_base_text += f"• Número de Plantios: {len(st.session_state.get('plantios', {}))}\n"
        dados_base_text += f"• Número de Despesas: {len(st.session_state.get('despesas', []))}\n"
        dados_base_text += f"• Número de Empréstimos: {len(st.session_state.get('emprestimos', []))}\n\n"
        
        dados_base_text += "PARÂMETROS DE CENÁRIO:\n"
        dados_base_text += f"• Receita Pessimista: -{st.session_state.get('pess_receita', 15)}%\n"
        dados_base_text += f"• Despesa Pessimista: +{st.session_state.get('pess_despesas', 10)}%\n"
        dados_base_text += f"• Receita Otimista: +{st.session_state.get('otm_receita', 10)}%\n"
        dados_base_text += f"• Despesa Otimista: -{st.session_state.get('otm_despesas', 10)}%\n\n"
        
        dados_base_text += "INFLAÇÃO PROJETADA:\n"
        for i, ano in enumerate(anos):
            inflacao = st.session_state.get('inflacoes', [4.0] * len(anos))[i]
            dados_base_text += f"• {ano}: {inflacao:.1f}%\n"
        
        content.text = dados_base_text
        
        # SEÇÃO CENÁRIO PROJETADO
        
        # SLIDE 4: DRE Completo - Projetado (TABELA)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "📊 DRE Consolidado - Projetado"
        
        dre_data = all_dre_data["Projetado"]
        dre_dict = {
            "Receita Operacional": dre_data["Receita"],
            "(-) Impostos s/ Venda": [-x for x in dre_data["Impostos Sobre Venda"]],
            "(-) Despesas Operacionais": [-x for x in dre_data["Despesas Operacionais"]],
            "(-) Despesas Administrativas": [-x for x in dre_data["Despesas Administrativas"]],
            "(-) Despesas RH": [-x for x in dre_data["Despesas RH"]],
            "(-) Despesas Extra Op.": [-x for x in dre_data["Despesas Extra Operacional"]],
            "(-) Dividendos": [-x for x in dre_data["Dividendos"]],
            "(-) Impostos s/ Resultado": [-x for x in dre_data["Impostos Sobre Resultado"]],
            "(=) LUCRO LÍQUIDO": dre_data["Lucro Líquido"]
        }
        
        criar_tabela_slide("📊 DRE Consolidado - Projetado", dre_dict, anos)

        # SLIDE 5: INDICADORES - PROJETADO (TABELA)
        criar_tabela_indicadores("📊 Indicadores Consolidado - Projetado", all_indicators["Projetado"], anos)

        # SLIDE 6: PARECER - PROJETADO (texto formatado)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "📊 Parecer Consolidado - Projetado"
        
        indicators = all_indicators["Projetado"]
        margem_media = np.mean(indicators["Margem Líquida (%)"])
        retorno_medio = np.mean(indicators["Retorno por Real Gasto"])
        liquidez_media = np.mean(indicators["Liquidez Operacional"])
        
        parecer_text = "PARECER FINANCEIRO - PROJETADO:\n\n"
        
        if margem_media < 10:
            parecer_text += f"🔴 MARGEM LÍQUIDA BAIXA ({margem_media:.1f}%):\n"
            parecer_text += "• Rentabilidade abaixo do ideal\n"
            parecer_text += "• Revisar estrutura de custos\n"
            parecer_text += "• Renegociar preços de venda\n\n"
        else:
            parecer_text += f"✅ MARGEM LÍQUIDA SAUDÁVEL ({margem_media:.1f}%):\n"
            parecer_text += "• Boa rentabilidade operacional\n"
            parecer_text += "• Manter eficiência atual\n\n"
        
        if retorno_medio < 0.2:
            parecer_text += f"🔴 BAIXO RETORNO ({retorno_medio:.2f}):\n"
            parecer_text += "• Investimentos com baixo retorno\n"
            parecer_text += "• Otimizar aplicação de recursos\n\n"
        else:
            parecer_text += f"✅ RETORNO ADEQUADO ({retorno_medio:.2f}):\n"
            parecer_text += "• Boa eficiência dos investimentos\n\n"
        
        if liquidez_media < 1.5:
            parecer_text += f"⚠️ LIQUIDEZ BAIXA ({liquidez_media:.2f}):\n"
            parecer_text += "• Risco de fluxo de caixa\n"
            parecer_text += "• Buscar linhas de crédito\n\n"
        else:
            parecer_text += f"✅ LIQUIDEZ CONFORTÁVEL ({liquidez_media:.2f}):\n"
            parecer_text += "• Boa capacidade operacional\n\n"
        
        if indicators['CAGR Lucro Líquido (%)'] < 0:
            parecer_text += f"🔴 CRESCIMENTO NEGATIVO ({indicators['CAGR Lucro Líquido (%)']:.1f}%):\n"
            parecer_text += "• Revisar estratégia\n"
            parecer_text += "• Analisar custos e preços"
        else:
            parecer_text += f"✅ CRESCIMENTO POSITIVO ({indicators['CAGR Lucro Líquido (%)']:.1f}%):\n"
            parecer_text += "• Trajetória sustentável"
        
        content.text = parecer_text

        # Por cultura - Projetado (se disponível)
        if all_indicators_cultura_cenarios and "Projetado" in all_indicators_cultura_cenarios:
            for cultura in all_indicators_cultura_cenarios["Projetado"].keys():
                indicators_cultura = all_indicators_cultura_cenarios["Projetado"][cultura]
                
                # DRE por cultura (TABELA ESTIMADA)
                receitas_por_cultura = st.session_state.get('receitas_por_cultura_cenarios', {}).get('Projetado', {})
                
                if cultura in receitas_por_cultura:
                    receitas_cultura = [receitas_por_cultura[cultura].get(ano, 0) for ano in anos]
                    margem_cultura = np.mean(indicators_cultura["Margem Líquida (%)"])
                    
                    custos_estimados = []
                    lucros_estimados = []
                    
                    for receita_ano in receitas_cultura:
                        lucro_ano = receita_ano * (margem_cultura / 100)
                        custo_ano = receita_ano - lucro_ano
                        custos_estimados.append(-custo_ano)  # Negativo para mostrar como despesa
                        lucros_estimados.append(lucro_ano)
                    
                    dre_cultura_dict = {
                        "(+) Receita": receitas_cultura,
                        "(-) Custos e Despesas": custos_estimados,
                        "(=) Lucro Líquido": lucros_estimados
                    }
                    
                    criar_tabela_slide(f"📊 DRE - {cultura} - Projetado", dre_cultura_dict, anos)
                    
                    # Indicadores por cultura (TABELA)
                    criar_tabela_indicadores(f"📊 Indicadores - {cultura} - Projetado", indicators_cultura, anos)
                    
                    # Parecer por cultura (TEXTO)
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = f"📊 Parecer - {cultura} - Projetado"
                    
                    margem_media_cultura = np.mean(indicators_cultura["Margem Líquida (%)"])
                    retorno_medio_cultura = np.mean(indicators_cultura["Retorno por Real Gasto"])
                    
                    parecer_cultura_text = f"PARECER - {cultura.upper()}:\n\n"
                    
                    if margem_media_cultura < 10:
                        parecer_cultura_text += f"🔴 MARGEM BAIXA ({margem_media_cultura:.1f}%)\n"
                        parecer_cultura_text += "• Otimizar técnicas de cultivo\n"
                        parecer_cultura_text += "• Revisar custos de insumos\n\n"
                    else:
                        parecer_cultura_text += f"✅ MARGEM SAUDÁVEL ({margem_media_cultura:.1f}%)\n"
                        parecer_cultura_text += "• Manter práticas atuais\n\n"
                    
                    if retorno_medio_cultura < 0.2:
                        parecer_cultura_text += f"🔴 BAIXO RETORNO ({retorno_medio_cultura:.2f})\n"
                        parecer_cultura_text += "• Revisar investimentos\n\n"
                    else:
                        parecer_cultura_text += f"✅ RETORNO ADEQUADO ({retorno_medio_cultura:.2f})\n\n"
                    
                    if indicators_cultura['CAGR Lucro Líquido (%)'] < 0:
                        parecer_cultura_text += f"🔴 CRESCIMENTO NEGATIVO\n"
                        parecer_cultura_text += "• Reavaliar viabilidade"
                    else:
                        parecer_cultura_text += f"✅ CRESCIMENTO POSITIVO\n"
                        parecer_cultura_text += "• Cultura sustentável"
                    
                    content.text = parecer_cultura_text

        # SEÇÃO CENÁRIO PESSIMISTA
        
        # DRE, Indicadores e Parecer - Pessimista (consolidado + culturas)
        dre_data_pess = all_dre_data["Pessimista"]
        dre_dict_pess = {
            "Receita Operacional": dre_data_pess["Receita"],
            "(-) Impostos s/ Venda": [-x for x in dre_data_pess["Impostos Sobre Venda"]],
            "(-) Despesas Operacionais": [-x for x in dre_data_pess["Despesas Operacionais"]],
            "(-) Despesas Administrativas": [-x for x in dre_data_pess["Despesas Administrativas"]],
            "(-) Despesas RH": [-x for x in dre_data_pess["Despesas RH"]],
            "(-) Despesas Extra Op.": [-x for x in dre_data_pess["Despesas Extra Operacional"]],
            "(-) Dividendos": [-x for x in dre_data_pess["Dividendos"]],
            "(-) Impostos s/ Resultado": [-x for x in dre_data_pess["Impostos Sobre Resultado"]],
            "(=) LUCRO LÍQUIDO": dre_data_pess["Lucro Líquido"]
        }
        
        criar_tabela_slide("📉 DRE Consolidado - Pessimista", dre_dict_pess, anos)
        criar_tabela_indicadores("📉 Indicadores Consolidado - Pessimista", all_indicators["Pessimista"], anos)
        
        # Por cultura - Pessimista
        if all_indicators_cultura_cenarios and "Pessimista" in all_indicators_cultura_cenarios:
            for cultura in all_indicators_cultura_cenarios["Pessimista"].keys():
                # Simular DRE por cultura (baseado nos indicadores)
                indicators_cultura = all_indicators_cultura_cenarios["Pessimista"][cultura]
                
                # Criar slide DRE estimado para cultura
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"📊 DRE - {cultura} - Pessimista"
                
                # Estimar valores baseados nas receitas por cultura
                receitas_por_cultura = st.session_state.get('receitas_por_cultura_cenarios', {}).get('Pessimista', {})
                
                if cultura in receitas_por_cultura:
                    receita_cultura_total = sum(receitas_por_cultura[cultura].get(ano, 0) for ano in anos)
                    margem_cultura = np.mean(indicators_cultura["Margem Líquida (%)"])
                    lucro_estimado = receita_cultura_total * (margem_cultura / 100)
                    
                    dre_cultura_text = f"DRE ESTIMADO - {cultura.upper()}:\n\n"
                    dre_cultura_text += f"(+) Receita Total (5 anos): R$ {receita_cultura_total:,.0f}\n"
                    dre_cultura_text += f"(-) Custos e Despesas: R$ {receita_cultura_total - lucro_estimado:,.0f}\n"
                    dre_cultura_text += f"(=) Lucro Líquido: R$ {lucro_estimado:,.0f}\n\n"
                    dre_cultura_text += f"MARGEM LÍQUIDA: {margem_cultura:.1f}%\n\n"
                    
                    hectares_cultura = sum(
                        plantio.get('hectares', 0) 
                        for plantio in st.session_state.get('plantios', {}).values() 
                        if plantio.get('cultura') == cultura
                    )
                    
                    if hectares_cultura > 0:
                        receita_por_ha = receita_cultura_total / hectares_cultura
                        lucro_por_ha = lucro_estimado / hectares_cultura
                        
                        dre_cultura_text += f"ANÁLISE POR HECTARE:\n"
                        dre_cultura_text += f"• Receita/Ha: R$ {receita_por_ha:,.0f}\n"
                        dre_cultura_text += f"• Lucro/Ha: R$ {lucro_por_ha:,.0f}\n"
                        dre_cultura_text += f"• Área Total: {hectares_cultura:.1f} hectares"
                    
                    content.text = dre_cultura_text
                else:
                    content.text = f"DRE para {cultura} não disponível - dados insuficientes"
                
                # Indicadores por cultura - Pessimista
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"📊 Indicadores - {cultura} - Pessimista"
                
                margem_media_cultura = np.mean(indicators_cultura["Margem Líquida (%)"])
                retorno_medio_cultura = np.mean(indicators_cultura["Retorno por Real Gasto"])
                liquidez_media_cultura = np.mean(indicators_cultura["Liquidez Operacional"])
                roa_medio_cultura = np.mean(indicators_cultura["ROA (%)"])
                
                indicadores_cultura_text = f"INDICADORES - {cultura.upper()}:\n\n"
                indicadores_cultura_text += "RENTABILIDADE:\n"
                indicadores_cultura_text += f"• Margem Líquida Média: {margem_media_cultura:.2f}%\n"
                indicadores_cultura_text += f"• Retorno por Real Gasto: {retorno_medio_cultura:.2f}\n"
                indicadores_cultura_text += f"• ROA Médio: {roa_medio_cultura:.2f}%\n\n"
                
                indicadores_cultura_text += "EFICIÊNCIA:\n"
                indicadores_cultura_text += f"• Liquidez Operacional: {liquidez_media_cultura:.2f}\n"
                
                indicadores_cultura_text += f"\nCRESCIMENTO:\n"
                indicadores_cultura_text += f"• CAGR Receita: {indicators_cultura['CAGR Receita (%)']:.2f}%\n"
                indicadores_cultura_text += f"• CAGR Lucro: {indicators_cultura['CAGR Lucro Líquido (%)']:.2f}%\n"
                
                content.text = indicadores_cultura_text
                
                # Parecer por cultura - Pessimista
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"📊 Parecer - {cultura} - Pessimista"
                
                parecer_cultura_text = f"PARECER - {cultura.upper()}:\n\n"
                
                if margem_media_cultura < 10:
                    parecer_cultura_text += f"🔴 MARGEM BAIXA ({margem_media_cultura:.1f}%)\n"
                    parecer_cultura_text += "• Otimizar técnicas de cultivo\n"
                    parecer_cultura_text += "• Revisar custos de insumos\n\n"
                else:
                    parecer_cultura_text += f"✅ MARGEM SAUDÁVEL ({margem_media_cultura:.1f}%)\n"
                    parecer_cultura_text += "• Manter práticas atuais\n\n"
                
                if retorno_medio_cultura < 0.2:
                    parecer_cultura_text += f"🔴 BAIXO RETORNO ({retorno_medio_cultura:.2f})\n"
                    parecer_cultura_text += "• Revisar investimentos\n\n"
                else:
                    parecer_cultura_text += f"✅ RETORNO ADEQUADO ({retorno_medio_cultura:.2f})\n\n"
                
                if indicators_cultura['CAGR Lucro Líquido (%)'] < 0:
                    parecer_cultura_text += f"🔴 CRESCIMENTO NEGATIVO\n"
                    parecer_cultura_text += "• Reavaliar viabilidade\n"
                else:
                    parecer_cultura_text += f"✅ CRESCIMENTO POSITIVO\n"
                    parecer_cultura_text += "• Cultura sustentável\n"
                
                content.text = parecer_cultura_text
        
        # SEÇÃO CENÁRIO OTIMISTA
        
        # DRE, Indicadores e Parecer - Otimista (consolidado + culturas)  
        dre_data_otm = all_dre_data["Otimista"]
        dre_dict_otm = {
            "Receita Operacional": dre_data_otm["Receita"],
            "(-) Impostos s/ Venda": [-x for x in dre_data_otm["Impostos Sobre Venda"]],
            "(-) Despesas Operacionais": [-x for x in dre_data_otm["Despesas Operacionais"]],
            "(-) Despesas Administrativas": [-x for x in dre_data_otm["Despesas Administrativas"]],
            "(-) Despesas RH": [-x for x in dre_data_otm["Despesas RH"]],
            "(-) Despesas Extra Op.": [-x for x in dre_data_otm["Despesas Extra Operacional"]],
            "(-) Dividendos": [-x for x in dre_data_otm["Dividendos"]],
            "(-) Impostos s/ Resultado": [-x for x in dre_data_otm["Impostos Sobre Resultado"]],
            "(=) LUCRO LÍQUIDO": dre_data_otm["Lucro Líquido"]
        }
        
        criar_tabela_slide("📈 DRE Consolidado - Otimista", dre_dict_otm, anos)
        criar_tabela_indicadores("📈 Indicadores Consolidado - Otimista", all_indicators["Otimista"], anos)

        # SLIDE: RECEITA vs LUCRO LÍQUIDO POR CENÁRIO (com gráfico)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "📈 Receita vs Lucro Líquido por Cenário"
        
        try:
            # Criar gráfico comparativo
            fig, ax = plt.subplots(figsize=(12, 7))
            x = range(len(anos))
            width = 0.25
            colors_map = {'Projetado': '#1f77b4', 'Pessimista': '#ff7f0e', 'Otimista': '#2ca02c'}
            
            for i, cenario in enumerate(nomes_cenarios):
                receitas = all_dre_data[cenario]["Receita"]
                lucros = all_dre_data[cenario]["Lucro Líquido"]
                
                ax.bar([p + width*i for p in x], receitas, width, 
                       label=f'Receita ({cenario})', color=colors_map[cenario], alpha=0.7)
                ax.bar([p + width*i for p in x], lucros, width,
                       label=f'Lucro ({cenario})', color=colors_map[cenario], alpha=0.4)
            
            ax.set_xlabel('Anos', fontsize=12)
            ax.set_ylabel('Valores (R$)', fontsize=12)
            ax.set_xticks([p + width for p in x])
            ax.set_xticklabels(anos)
            ax.legend()
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'R$ {x/1e6:.1f}M'))
            plt.tight_layout()
            
            # Salvar gráfico
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            plt.close()
            
            # Adicionar imagem ao slide (redimensionar para caber com texto)
            slide.shapes.add_picture(img_buffer, Inches(1), Inches(2), Inches(8), Inches(4))
            
            # Adicionar análise textual abaixo do gráfico
            grafico_text = "ANÁLISE COMPARATIVA:\n\n"
            
            lucro_proj = sum(all_dre_data["Projetado"]["Lucro Líquido"])
            lucro_pess = sum(all_dre_data["Pessimista"]["Lucro Líquido"])
            lucro_otm = sum(all_dre_data["Otimista"]["Lucro Líquido"])
            
            diff_pess = ((lucro_pess - lucro_proj) / lucro_proj * 100) if lucro_proj != 0 else 0
            diff_otm = ((lucro_otm - lucro_proj) / lucro_proj * 100) if lucro_proj != 0 else 0
            
            grafico_text += f"• Cenário Pessimista: {diff_pess:+.1f}% vs Projetado\n"
            grafico_text += f"• Cenário Otimista: {diff_otm:+.1f}% vs Projetado"
            
            # Adicionar texto acima do gráfico
            text_shape = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.6))
            text_frame = text_shape.text_frame
            text_frame.text = grafico_text
            
        except Exception as chart_error:
            # Fallback se não conseguir criar o gráfico
            content.text = f"COMPARATIVO DE RECEITA vs LUCRO LÍQUIDO:\n\nDados por cenário (5 anos):\n\n" + \
                          f"📊 PROJETADO:\n• Receita: R$ {sum(all_dre_data['Projetado']['Receita']):,.0f}\n• Lucro: R$ {sum(all_dre_data['Projetado']['Lucro Líquido']):,.0f}\n\n" + \
                          f"📉 PESSIMISTA:\n• Receita: R$ {sum(all_dre_data['Pessimista']['Receita']):,.0f}\n• Lucro: R$ {sum(all_dre_data['Pessimista']['Lucro Líquido']):,.0f}\n\n" + \
                          f"📈 OTIMISTA:\n• Receita: R$ {sum(all_dre_data['Otimista']['Receita']):,.0f}\n• Lucro: R$ {sum(all_dre_data['Otimista']['Lucro Líquido']):,.0f}"
        
        # SLIDE: AGRADECIMENTO
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Obrigado!"
        subtitle.text = f"Relatório gerado pelo Sistema Gestor de Plantio\n{datetime.now().strftime('%d/%m/%Y às %H:%M')}\n\nAnálise Completa Finalizada"
        
        # Salvar apresentação
        output_ppt = BytesIO()
        prs.save(output_ppt)
        output_ppt.seek(0)
        
        return output_ppt
        
    except ImportError:
        return None
    except Exception as e:
        st.error(f"Erro ao gerar PowerPoint: {str(e)}")
        return None