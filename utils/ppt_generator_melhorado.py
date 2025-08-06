import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from datetime import datetime

def criar_relatorio_ppt_completo(all_indicators, all_dre_data, df_culturas_for_excel, nomes_cenarios, anos, all_indicators_cultura_cenarios=None):
    """
    Gera um relatório COMPLETO em PowerPoint com TODAS as informações do 5_indicadores.py
    Inclui: DREs, Indicadores, Pareceres, Fluxos de Caixa, Análises por Cultura, etc.
    """
    try:
        # Importações necessárias
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Criar nova apresentação
        prs = Presentation()
        st.write("🔧 Iniciando geração PowerPoint COMPLETO baseado em 5_indicadores.py...")
        
        def criar_slide_texto(titulo, texto):
            """Cria slide apenas com texto formatado"""
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = titulo
            
            # Configurar o conteúdo de texto
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            # Dividir texto em parágrafos
            paragraphs = texto.split('\n')
            for i, paragraph in enumerate(paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = paragraph
                p.font.size = Pt(11)
                if paragraph.startswith('• **') or paragraph.startswith('🔴') or paragraph.startswith('⚠️') or paragraph.startswith('✅'):
                    p.font.bold = True
            
            return slide

        def criar_slide_com_tabela(title_text, df, format_financeiro=True):
            """Cria slide com DataFrame formatado como tabela"""
            if df.empty:
                return criar_slide_texto(title_text, "Dados não disponíveis para esta análise.")
                
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = title_text
            
            # Remover placeholder de conteúdo
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                sp = content_placeholder.element
                sp.getparent().remove(sp)
            
            # Criar tabela
            rows = len(df) + 1
            cols = len(df.columns)
            
            left = Inches(0.2)
            top = Inches(1.5)
            width = Inches(9.6)
            height = Inches(5.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Cabeçalho
            for j, col_name in enumerate(df.columns):
                cell = table.cell(0, j)
                cell.text = str(col_name)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Dados
            for i, (idx, row) in enumerate(df.iterrows()):
                for j, value in enumerate(row):
                    cell = table.cell(i + 1, j)
                    
                    if pd.isna(value):
                        cell.text = ""
                    elif isinstance(value, (int, float)):
                        if format_financeiro:
                            if "%" in str(df.columns[j]) or "Margem" in str(df.columns[j]):
                                cell.text = f"{value:.2f}%"
                            elif "R$" in str(df.columns[j]) or abs(value) >= 1000:
                                cell.text = f"R$ {value:,.0f}"
                            elif "CAGR" in str(df.columns[j]):
                                cell.text = f"{value:.2f}%"
                            else:
                                cell.text = f"{value:.2f}"
                        else:
                            cell.text = str(value)
                    else:
                        cell.text = str(value)
                    
                    cell.text_frame.paragraphs[0].font.size = Pt(8)
            
            return slide

        def gerar_parecer_consolidado_detalhado(indicators_cenario, cenario_name):
            """Gera parecer financeiro detalhado idêntico ao do 5_indicadores.py"""
            
            margem_media = np.mean(indicators_cenario["Margem Líquida (%)"])
            retorno_medio = np.mean(indicators_cenario["Retorno por Real Gasto"])
            liquidez_media = np.mean(indicators_cenario["Liquidez Operacional"])
            endividamento_medio = np.mean(indicators_cenario["Endividamento (%)"])
            produtividade_media = np.mean(indicators_cenario["Produtividade por Hectare (R$/ha)"])
            custo_receita_media = np.mean(indicators_cenario["Custo por Receita (%)"])
            
            # Tratar DSCR que pode ter valores infinitos
            dscr_values = [x for x in indicators_cenario["DSCR"] if x != float("inf")]
            dscr_medio = np.mean(dscr_values) if dscr_values else float("inf")
            
            break_even_media = np.mean(indicators_cenario["Break-Even Yield (sacas/ha)"])
            roa_medio = np.mean(indicators_cenario["ROA (%)"])
            
            parecer_items = []
            
            # Cabeçalho
            parecer_items.extend([
                f"ANÁLISE FINANCEIRA DETALHADA - CENÁRIO {cenario_name.upper()}",
                "=" * 60,
                "",
                "📊 INDICADORES DE RENTABILIDADE:"
            ])
            
            # Margem Líquida
            if margem_media < 10:
                parecer_items.extend([
                    f"🔴 MARGEM LÍQUIDA BAIXA ({margem_media:.2f}%)",
                    "   • Rentabilidade abaixo do ideal para o agronegócio",
                    "   • Considere renegociar preços com fornecedores",
                    "   • Investir em culturas de maior valor agregado",
                    "   • Revisar técnicas de plantio e colheita"
                ])
            elif margem_media < 20:
                parecer_items.extend([
                    f"⚠️ MARGEM LÍQUIDA MODERADA ({margem_media:.2f}%)",
                    "   • Rentabilidade aceitável, mas pode melhorar",
                    "   • Monitore custos para manter consistência",
                    "   • Busque oportunidades de otimização"
                ])
            else:
                parecer_items.extend([
                    f"✅ MARGEM LÍQUIDA EXCELENTE ({margem_media:.2f}%)",
                    "   • Excelente rentabilidade para o setor",
                    "   • Mantenha as práticas atuais",
                    "   • Considere reinvestir em expansão"
                ])
            
            parecer_items.append("")
            
            # Retorno por Real Gasto
            parecer_items.append("💰 EFICIÊNCIA DOS INVESTIMENTOS:")
            if retorno_medio < 0.2:
                parecer_items.extend([
                    f"🔴 BAIXO RETORNO POR REAL GASTO ({retorno_medio:.2f})",
                    "   • Cada R$ 1,00 gasto retorna menos de R$ 0,20",
                    "   • Gastos com baixo retorno financeiro",
                    "   • Avalie redução de despesas operacionais",
                    "   • Otimize processos agrícolas e logísticos"
                ])
            else:
                parecer_items.extend([
                    f"✅ RETORNO ADEQUADO POR REAL GASTO ({retorno_medio:.2f})",
                    f"   • Cada R$ 1,00 gasto retorna R$ {retorno_medio:.2f}",
                    "   • Investimentos geram retorno satisfatório",
                    "   • Considere reinvestir em tecnologia"
                ])
            
            parecer_items.append("")
            
            # Liquidez Operacional
            parecer_items.append("🏦 LIQUIDEZ E CAPACIDADE OPERACIONAL:")
            if liquidez_media < 1.5:
                parecer_items.extend([
                    f"🔴 LIQUIDEZ OPERACIONAL BAIXA ({liquidez_media:.2f})",
                    "   • Risco de dificuldades para cobrir custos",
                    "   • Negocie prazos de pagamento com fornecedores",
                    "   • Busque linhas de crédito de curto prazo",
                    "   • Monitore fluxo de caixa diariamente"
                ])
            else:
                parecer_items.extend([
                    f"✅ LIQUIDEZ OPERACIONAL CONFORTÁVEL ({liquidez_media:.2f})",
                    "   • Boa capacidade de sustentar operações",
                    "   • Mantenha reservas para safras incertas",
                    "   • Considere investimentos estratégicos"
                ])
            
            parecer_items.append("")
            
            # Endividamento
            parecer_items.append("📈 ESTRUTURA DE ENDIVIDAMENTO:")
            if endividamento_medio > 30:
                parecer_items.extend([
                    f"🔴 ALTO ENDIVIDAMENTO ({endividamento_medio:.2f}%)",
                    "   • Dívidas elevadas comprometem rentabilidade",
                    "   • Priorize quitação de empréstimos de alto custo",
                    "   • Renegocie taxas de juros com instituições",
                    "   • Evite novos financiamentos não essenciais"
                ])
            else:
                parecer_items.extend([
                    f"✅ ENDIVIDAMENTO CONTROLADO ({endividamento_medio:.2f}%)",
                    "   • Dívidas em nível gerenciável",
                    "   • Considere investimentos estratégicos",
                    "   • Expansão de área plantada pode ser viável"
                ])
            
            parecer_items.append("")
            
            # Produtividade
            parecer_items.append("🌱 PRODUTIVIDADE E EFICIÊNCIA:")
            produtividade_status = "Boa produtividade" if produtividade_media > 5000 else "Produtividade pode ser melhorada"
            parecer_items.extend([
                f"📊 PRODUTIVIDADE POR HECTARE: R$ {produtividade_media:,.0f}/ha",
                f"   • {produtividade_status}",
                "   • Compare com benchmarks da região",
                "   • Avalie técnicas de manejo e insumos"
            ])
            
            parecer_items.append("")
            
            # Custo por Receita
            parecer_items.append("💸 GESTÃO DE CUSTOS:")
            if custo_receita_media > 70:
                parecer_items.extend([
                    f"🔴 CUSTO POR RECEITA ALTO ({custo_receita_media:.2f}%)",
                    "   • Custos operacionais consomem grande parte da receita",
                    "   • Analise insumos e processos para reduzir despesas",
                    "   • Negocie melhores preços com fornecedores",
                    "   • Otimize logística e armazenamento"
                ])
            else:
                parecer_items.extend([
                    f"✅ CUSTO POR RECEITA CONTROLADO ({custo_receita_media:.2f}%)",
                    "   • Boa gestão de custos operacionais",
                    "   • Continue monitorando preços de insumos",
                    "   • Mantenha controle rigoroso dos gastos"
                ])
            
            parecer_items.append("")
            
            # DSCR
            parecer_items.append("🏛️ CAPACIDADE DE PAGAMENTO DE DÍVIDAS:")
            if dscr_medio != float("inf") and dscr_medio < 1.25:
                parecer_items.extend([
                    f"🔴 DSCR BAIXO ({dscr_medio:.2f})",
                    "   • Risco de dificuldades no pagamento de dívidas",
                    "   • Considere reestruturar financiamentos",
                    "   • Aumente receita ou reduza despesas",
                    "   • Monitore fluxo de caixa mensalmente"
                ])
            else:
                dscr_text = f"{dscr_medio:.2f}" if dscr_medio != float("inf") else "∞"
                parecer_items.extend([
                    f"✅ DSCR ADEQUADO ({dscr_text})",
                    "   • Boa capacidade de cobrir dívidas",
                    "   • Mantenha lucro operacional estável",
                    "   • Considere expansão controlada"
                ])
            
            parecer_items.append("")
            
            # Break-Even Yield
            parecer_items.append("⚖️ PONTO DE EQUILÍBRIO:")
            risco_safra = "Alto risco em safras ruins" if break_even_media > 50 else "Risco moderado em cenários adversos"
            parecer_items.extend([
                f"📊 BREAK-EVEN YIELD: {break_even_media:.1f} sacas/ha",
                f"   • {risco_safra}",
                "   • Diversifique culturas para reduzir risco",
                "   • Invista em seguro agrícola"
            ])
            
            parecer_items.append("")
            
            # ROA
            parecer_items.append("🏭 RETORNO SOBRE ATIVOS:")
            if roa_medio < 5:
                parecer_items.extend([
                    f"🔴 ROA BAIXO ({roa_medio:.2f}%)",
                    "   • Baixa eficiência no uso de ativos",
                    "   • Avalie venda de ativos ociosos",
                    "   • Invista em equipamentos mais produtivos",
                    "   • Otimize uso da terra disponível"
                ])
            else:
                parecer_items.extend([
                    f"✅ ROA ADEQUADO ({roa_medio:.2f}%)",
                    "   • Boa utilização dos ativos",
                    "   • Considere expansão controlada",
                    "   • Mantenha eficiência operacional"
                ])
            
            parecer_items.append("")
            
            # CAGR
            parecer_items.append("📈 ANÁLISE DE CRESCIMENTO:")
            cagr_receita = indicators_cenario.get("CAGR Receita (%)", 0)
            cagr_lucro = indicators_cenario.get("CAGR Lucro Líquido (%)", 0)
            
            parecer_items.extend([
                f"📊 CAGR RECEITA (5 anos): {cagr_receita:.2f}%/ano",
                f"📊 CAGR LUCRO LÍQUIDO (5 anos): {cagr_lucro:.2f}%/ano"
            ])
            
            if cagr_lucro < 0:
                parecer_items.extend([
                    "🔴 CRESCIMENTO NEGATIVO DO LUCRO",
                    "   • Lucro em queda ao longo dos anos",
                    "   • Revisar estratégias de custo, preço e produtividade",
                    "   • Considerar mudança no mix de culturas"
                ])
            elif cagr_lucro < 5:
                parecer_items.extend([
                    "⚠️ CRESCIMENTO BAIXO DO LUCRO",
                    "   • Crescimento abaixo da média do setor",
                    "   • Buscar oportunidades de melhoria"
                ])
            else:
                parecer_items.extend([
                    "✅ CRESCIMENTO SAUDÁVEL DO LUCRO",
                    "   • Lucro em trajetória positiva consistente",
                    "   • Considere reinvestir em áreas estratégicas"
                ])
            
            return "\n".join(parecer_items)

        def gerar_parecer_cultura_detalhado(indicators_cultura, cultura, hectares_cultura, cenario):
            """Gera parecer detalhado por cultura"""
            if not indicators_cultura:
                return f"Dados insuficientes para análise da cultura {cultura}"
            
            margem_cultura = np.mean(indicators_cultura.get("Margem Líquida (%)", [0]))
            retorno_cultura = np.mean(indicators_cultura.get("Retorno por Real Gasto", [0]))
            liquidez_cultura = np.mean(indicators_cultura.get("Liquidez Operacional", [0]))
            roa_cultura = np.mean(indicators_cultura.get("ROA (%)", [0]))
            custo_receita_cultura = np.mean(indicators_cultura.get("Custo por Receita (%)", [0]))
            
            # Calcular receita total da cultura
            receitas_por_cultura = st.session_state.get('receitas_por_cultura_cenarios', {}).get(cenario, {})
            receita_total_cultura = sum(receitas_por_cultura.get(cultura, {}).get(str(ano), 0) for ano in anos) if cultura in receitas_por_cultura else 0
            lucro_total_cultura = receita_total_cultura * (margem_cultura / 100) if margem_cultura > 0 else 0
            
            parecer_items = [
                f"ANÁLISE DETALHADA - CULTURA {cultura.upper()}",
                f"CENÁRIO: {cenario.upper()}",
                "=" * 60,
                "",
                "📊 DADOS OPERACIONAIS:",
                f"   • Área cultivada: {hectares_cultura:.1f} hectares",
                f"   • Receita total projetada ({len(anos)} anos): R$ {receita_total_cultura:,.0f}",
                f"   • Lucro estimado: R$ {lucro_total_cultura:,.0f}"
            ]
            
            if hectares_cultura > 0:
                receita_ha_ano = receita_total_cultura / (hectares_cultura * len(anos))
                lucro_ha_ano = lucro_total_cultura / (hectares_cultura * len(anos))
                parecer_items.extend([
                    f"   • Receita média por hectare/ano: R$ {receita_ha_ano:,.0f}",
                    f"   • Lucro médio por hectare/ano: R$ {lucro_ha_ano:,.0f}"
                ])
            
            parecer_items.extend([
                "",
                "💰 ANÁLISE DE RENTABILIDADE:"
            ])
            
            # Análise de performance
            if margem_cultura < 10:
                parecer_items.extend([
                    f"🔴 MARGEM LÍQUIDA BAIXA: {margem_cultura:.1f}%",
                    "   • Ação urgente necessária",
                    "   • Revisar custos de insumos específicos",
                    "   • Avaliar técnicas de manejo",
                    "   • Considerar variedades mais rentáveis"
                ])
            elif margem_cultura < 20:
                parecer_items.extend([
                    f"⚠️ MARGEM LÍQUIDA MODERADA: {margem_cultura:.1f}%",
                    "   • Performance aceitável",
                    "   • Buscar melhorias incrementais",
                    "   • Monitorar preços de mercado"
                ])
            else:
                parecer_items.extend([
                    f"✅ MARGEM LÍQUIDA EXCELENTE: {margem_cultura:.1f}%",
                    "   • Cultura muito rentável",
                    "   • Considerar expansão da área",
                    "   • Mantenha práticas atuais"
                ])
            
            # Eficiência
            parecer_items.extend([
                "",
                "⚡ INDICADORES DE EFICIÊNCIA:",
                f"   • Retorno por Real gasto: R$ {retorno_cultura:.2f}",
                f"   • Liquidez operacional: {liquidez_cultura:.2f}",
                f"   • ROA: {roa_cultura:.1f}%",
                f"   • Custo por receita: {custo_receita_cultura:.1f}%"
            ])
            
            # CAGR se disponível
            cagr_receita_cultura = indicators_cultura.get("CAGR Receita (%)", 0)
            cagr_lucro_cultura = indicators_cultura.get("CAGR Lucro Líquido (%)", 0)
            
            if cagr_lucro_cultura != 0:
                parecer_items.extend([
                    "",
                    "📈 CRESCIMENTO PROJETADO:",
                    f"   • CAGR Receita: {cagr_receita_cultura:.1f}%/ano",
                    f"   • CAGR Lucro: {cagr_lucro_cultura:.1f}%/ano"
                ])
                
                if cagr_lucro_cultura < 0:
                    parecer_items.append("   🔴 Tendência declinante - reavaliar viabilidade")
                elif cagr_lucro_cultura < 5:
                    parecer_items.append("   ⚠️ Crescimento lento - buscar otimizações")
                else:
                    parecer_items.append("   ✅ Crescimento saudável - manter estratégia")
            
            # Recomendação final
            parecer_items.extend([
                "",
                "🎯 RECOMENDAÇÃO ESTRATÉGICA:"
            ])
            
            if margem_cultura >= 15 and retorno_cultura >= 0.2 and cagr_lucro_cultura >= 0:
                parecer_items.extend([
                    "   ✅ CULTURA ALTAMENTE RECOMENDADA",
                    "   • Excelente performance em todos os indicadores",
                    "   • Priorizar aumento da área cultivada",
                    "   • Investir em tecnologia para esta cultura",
                    "   • Usar como referência para outras culturas"
                ])
            elif margem_cultura >= 10 and retorno_cultura >= 0.15:
                parecer_items.extend([
                    "   ✅ CULTURA RECOMENDADA",
                    "   • Performance adequada para manutenção",
                    "   • Manter área atual",
                    "   • Buscar melhorias graduais",
                    "   • Monitorar tendências de mercado"
                ])
            else:
                parecer_items.extend([
                    "   ⚠️ CULTURA REQUER ATENÇÃO ESPECIAL",
                    "   • Baixa performance financeira",
                    "   • Revisar completamente a estratégia",
                    "   • Considerar substituição por outras culturas",
                    "   • Avaliar custos de saída vs melhoria"
                ])
            
            return "\n".join(parecer_items)

        # ========== INÍCIO DA CRIAÇÃO DOS SLIDES ==========
        
        # SLIDE 1: CAPA PROFISSIONAL
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "📈 RELATÓRIO COMPLETO DE ANÁLISE FINANCEIRA"
        subtitle.text = f"""Gestão de Plantio - Análise Consolidada e por Cultura
        
Período de Análise: {anos[0]} - {anos[-1]}
Cenários: Projetado | Pessimista | Otimista

📊 Indicadores Financeiros Detalhados
🌱 Análise por Cultura e Cenário  
💰 Fluxos de Caixa Projetados
📋 Pareceres e Recomendações

Gerado em: {datetime.now().strftime('%d/%m/%Y às %H:%M')}
Base: Sistema 5_Indicadores.py"""
        
        # SLIDE 2: SUMÁRIO EXECUTIVO
        executivo_text = f"""RESUMO EXECUTIVO - ANÁLISE FINANCEIRA

📊 ESCOPO DA ANÁLISE:
• Período: {len(anos)} anos ({anos[0]} - {anos[-1]})
• Cenários analisados: {len(nomes_cenarios)} (Projetado, Pessimista, Otimista)
• Culturas avaliadas: {len(df_culturas_for_excel) if not df_culturas_for_excel.empty else 'Não especificado'}

💰 PRINCIPAIS MÉTRICAS (Cenário Projetado):
• Receita total projetada: R$ {sum(all_dre_data['Projetado']['Receita']):,.0f}
• Lucro líquido total: R$ {sum(all_dre_data['Projetado']['Lucro Líquido']):,.0f}
• Margem líquida média: {np.mean(all_indicators['Projetado']['Margem Líquida (%)']):,.1f}%
• ROA médio: {np.mean(all_indicators['Projetado']['ROA (%)']):,.1f}%

🎯 ESTRUTURA DO RELATÓRIO:
✓ Análise consolidada por cenário
✓ Indicadores financeiros detalhados  
✓ Pareceres técnicos especializados
✓ Análise individual por cultura
✓ Fluxos de caixa projetados
✓ Comparativos entre cenários
✓ Recomendações estratégicas"""
        
        criar_slide_texto("📋 Sumário Executivo", executivo_text)

        # SLIDE 3: EXPLICAÇÃO DOS INDICADORES
        explicacao_text = """ENTENDENDO OS INDICADORES FINANCEIROS

📊 INDICADORES DE RENTABILIDADE:
• Margem Líquida (%): Percentual do lucro em relação à receita
• Retorno por Real Gasto: Quanto cada R$ 1,00 gasto retorna de lucro
• ROA (%): Retorno sobre ativos - eficiência no uso de recursos

⚖️ INDICADORES DE LIQUIDEZ:
• Liquidez Operacional: Capacidade de cobrir custos operacionais
• DSCR: Capacidade de pagamento de dívidas

💸 INDICADORES DE CUSTOS:
• Custo por Receita (%): Proporção dos custos na receita
• Endividamento (%): Peso das dívidas na receita

🌱 INDICADORES AGRONEGÓCIO:
• Produtividade por Hectare (R$/ha): Receita gerada por hectare
• Break-Even Yield (sacas/ha): Produção mínima para cobrir custos
• Custo por Hectare (R$/ha): Custos totais por hectare

📈 INDICADORES DE CRESCIMENTO:
• CAGR Receita (%): Taxa de crescimento anual da receita
• CAGR Lucro (%): Taxa de crescimento anual do lucro"""
        
        criar_slide_texto("📚 Guia dos Indicadores Financeiros", explicacao_text)

        # ========== ANÁLISE CONSOLIDADA POR CENÁRIO ==========
        st.write("📊 Criando slides consolidados por cenário...")
        
        for i, cenario in enumerate(nomes_cenarios):
            emoji = "📊" if cenario == "Projetado" else "📉" if cenario == "Pessimista" else "📈"
            
            # DRE Consolidado
            if cenario in all_dre_data:
                dre_df = pd.DataFrame(all_dre_data[cenario])
                dre_df.index = [f"Ano {ano}" for ano in anos]
                criar_slide_com_tabela(f"{emoji} DRE Consolidado - {cenario}", dre_df)
            
            # Indicadores Consolidados
            if cenario in all_indicators:
                indicators_data = {}
                for key, values in all_indicators[cenario].items():
                    if isinstance(values, list) and len(values) == len(anos):
                        indicators_data[key] = values
                    elif isinstance(values, (int, float)):
                        # Para CAGR que são valores únicos
                        if key in ["CAGR Receita (%)", "CAGR Lucro Líquido (%)"]:
                            indicators_data[key] = [values] * len(anos)
                
                if indicators_data:
                    indicators_df = pd.DataFrame(indicators_data)
                    indicators_df.index = [f"Ano {ano}" for ano in anos]
                    
                    # Separar em duas tabelas se houver muitas colunas
                    if len(indicators_df.columns) > 6:
                        # Primeira tabela - indicadores principais
                        cols_principais = ["Margem Líquida (%)", "Retorno por Real Gasto", 
                                         "Liquidez Operacional", "ROA (%)", 
                                         "Produtividade por Hectare (R$/ha)", "Custo por Receita (%)"]
                        df_principais = indicators_df[[col for col in cols_principais if col in indicators_df.columns]]
                        criar_slide_com_tabela(f"{emoji} Indicadores Principais - {cenario}", df_principais)
                        
                        # Segunda tabela - indicadores complementares  
                        cols_complementares = [col for col in indicators_df.columns if col not in cols_principais]
                        if cols_complementares:
                            df_complementares = indicators_df[cols_complementares]
                            criar_slide_com_tabela(f"{emoji} Indicadores Complementares - {cenario}", df_complementares)
                    else:
                        criar_slide_com_tabela(f"{emoji} Indicadores Consolidados - {cenario}", indicators_df)
            
            # Parecer consolidado detalhado
            if cenario in all_indicators:
                parecer_text = gerar_parecer_consolidado_detalhado(all_indicators[cenario], cenario)
                criar_slide_texto(f"{emoji} Parecer Detalhado - {cenario}", parecer_text)

        # ========== COMPARATIVO ENTRE CENÁRIOS ==========
        st.write("📊 Criando comparativo entre cenários...")
        
        # Tabela comparativa final
        comparativo_data = []
        for cenario in nomes_cenarios:
            if cenario in all_indicators:
                receita_total = sum(all_dre_data[cenario]["Receita"])
                lucro_total = sum(all_dre_data[cenario]["Lucro Líquido"])
                margem_media = np.mean(all_indicators[cenario]["Margem Líquida (%)"])
                retorno_medio = np.mean(all_indicators[cenario]["Retorno por Real Gasto"])
                roa_medio = np.mean(all_indicators[cenario]["ROA (%)"])
                liquidez_media = np.mean(all_indicators[cenario]["Liquidez Operacional"])
                
                comparativo_data.append({
                    "Cenário": cenario,
                    f"Receita Total ({len(anos)}a)": receita_total,
                    f"Lucro Total ({len(anos)}a)": lucro_total,
                    "Margem Média (%)": margem_media,
                    "Retorno p/ Real": retorno_medio,
                    "ROA Médio (%)": roa_medio,
                    "Liquidez Média": liquidez_media
                })
        
        if comparativo_data:
            comparativo_df = pd.DataFrame(comparativo_data)
            criar_slide_com_tabela("📊 Resumo Comparativo - Todos os Cenários", comparativo_df)

        # Análise comparativa detalhada
        if len(comparativo_data) >= 2:
            lucro_proj = sum(all_dre_data['Projetado']['Lucro Líquido'])
            lucro_pess = sum(all_dre_data['Pessimista']['Lucro Líquido']) if 'Pessimista' in all_dre_data else lucro_proj
            lucro_otm = sum(all_dre_data['Otimista']['Lucro Líquido']) if 'Otimista' in all_dre_data else lucro_proj
            
            diferenca_pess = ((lucro_pess - lucro_proj) / lucro_proj * 100) if lucro_proj != 0 else 0
            diferenca_otm = ((lucro_otm - lucro_proj) / lucro_proj * 100) if lucro_proj != 0 else 0
            
            comparativo_text = f"""ANÁLISE COMPARATIVA ENTRE CENÁRIOS

📊 IMPACTO FINANCEIRO DOS CENÁRIOS:

💰 LUCRO TOTAL ({len(anos)} ANOS):
• Cenário Projetado: R$ {lucro_proj:,.0f}
• Cenário Pessimista: R$ {lucro_pess:,.0f} ({diferenca_pess:+.1f}%)
• Cenário Otimista: R$ {lucro_otm:,.0f} ({diferenca_otm:+.1f}%)

📈 VARIAÇÃO EM RELAÇÃO AO PROJETADO:
• Impacto Negativo (Pessimista): R$ {abs(lucro_pess - lucro_proj):,.0f}
• Potencial Positivo (Otimista): R$ {lucro_otm - lucro_proj:,.0f}
• Amplitude Total: R$ {abs(lucro_otm - lucro_pess):,.0f}

⚠️ ANÁLISE DE RISCO:
{'• Alto risco: variação pessimista > 20%' if abs(diferenca_pess) > 20 else '• Risco moderado: variação controlada'}
{'• Grande oportunidade: potencial otimista > 15%' if diferenca_otm > 15 else '• Oportunidade moderada'}

🎯 RECOMENDAÇÕES ESTRATÉGICAS:
• Preparar planos de contingência para cenário pessimista
• Identificar gatilhos para capturar oportunidades otimistas  
• Monitorar indicadores-chave para ajustes tempestivos
• Manter flexibilidade operacional e financeira"""
            
            criar_slide_texto("⚖️ Análise Comparativa Detalhada", comparativo_text)

        # ========== ANÁLISE POR CULTURA ==========
        if all_indicators_cultura_cenarios:
            st.write("🌱 Criando slides por cultura...")
            
            # Buscar dados complementares
            dre_por_cultura_cenarios = st.session_state.get('dre_por_cultura_cenarios', {})
            
            for cenario in nomes_cenarios:
                emoji = "📊" if cenario == "Projetado" else "📉" if cenario == "Pessimista" else "📈"
                
                if cenario in all_indicators_cultura_cenarios:
                    culturas = list(all_indicators_cultura_cenarios[cenario].keys())
                    
                    if culturas:
                        # Slide resumo das culturas no cenário
                        resumo_culturas = []
                        for cultura in culturas:
                            indicators_cultura = all_indicators_cultura_cenarios[cenario][cultura]
                            if indicators_cultura:
                                margem = np.mean(indicators_cultura.get("Margem Líquida (%)", [0]))
                                retorno = np.mean(indicators_cultura.get("Retorno por Real Gasto", [0]))
                                
                                # Dados operacionais
                                plantios_cultura = [p for p in st.session_state.get('plantios', {}).values() if p.get('cultura') == cultura]
                                hectares = sum(p.get('hectares', 0) for p in plantios_cultura)
                                
                                resumo_culturas.append({
                                    "Cultura": cultura,
                                    "Área (ha)": hectares,
                                    "Margem (%)": margem,
                                    "Retorno": retorno,
                                    "Status": "✅" if margem >= 15 else "⚠️" if margem >= 10 else "🔴"
                                })
                        
                        if resumo_culturas:
                            resumo_df = pd.DataFrame(resumo_culturas)
                            criar_slide_com_tabela(f"{emoji} Resumo por Cultura - {cenario}", resumo_df)
                    
                    # Para cada cultura individualmente
                    for cultura in culturas:
                        st.write(f"     Processando {cultura}...")
                        
                        # DRE da cultura
                        if cenario in dre_por_cultura_cenarios and cultura in dre_por_cultura_cenarios[cenario]:
                            dre_cultura = dre_por_cultura_cenarios[cenario][cultura]
                            if dre_cultura:
                                dre_cultura_df = pd.DataFrame(dre_cultura)
                                dre_cultura_df.index = [f"Ano {ano}" for ano in anos]
                                criar_slide_com_tabela(f"{emoji} DRE {cultura} - {cenario}", dre_cultura_df)
                        
                        # Indicadores da cultura
                        if cultura in all_indicators_cultura_cenarios[cenario]:
                            indicators_cultura = all_indicators_cultura_cenarios[cenario][cultura]
                            
                            indicators_cultura_data = {}
                            for key, values in indicators_cultura.items():
                                if isinstance(values, list) and len(values) == len(anos):
                                    indicators_cultura_data[key] = values
                                elif isinstance(values, (int, float)):
                                    if key in ["CAGR Receita (%)", "CAGR Lucro Líquido (%)"]:
                                        indicators_cultura_data[key] = [values] * len(anos)
                            
                            if indicators_cultura_data:
                                indicators_cultura_df = pd.DataFrame(indicators_cultura_data)
                                indicators_cultura_df.index = [f"Ano {ano}" for ano in anos]
                                criar_slide_com_tabela(f"{emoji} Indicadores {cultura} - {cenario}", indicators_cultura_df)
                            
                            # Parecer detalhado da cultura
                            hectares_cultura = sum(
                                p.get('hectares', 0) 
                                for p in st.session_state.get('plantios', {}).values() 
                                if p.get('cultura') == cultura
                            )
                            
                            parecer_cultura_text = gerar_parecer_cultura_detalhado(
                                indicators_cultura, cultura, hectares_cultura, cenario
                            )
                            criar_slide_texto(f"{emoji} Análise {cultura} - {cenario}", parecer_cultura_text)

        # ========== FLUXO DE CAIXA ==========
        st.write("💰 Criando slides de fluxo de caixa...")
        
        # Fluxo de caixa consolidado
        for cenario in nomes_cenarios:
            emoji = "📊" if cenario == "Projetado" else "📉" if cenario == "Pessimista" else "📈"
            
            if cenario in all_dre_data:
                dre_data = all_dre_data[cenario]
                
                fluxo_data = {
                    "Receita Operacional": dre_data["Receita"],
                    "(-) Impostos s/ Venda": [-x for x in dre_data["Impostos Sobre Venda"]],
                    "(-) Desp. Operacionais": [-x for x in dre_data["Despesas Operacionais"]],
                    "(-) Desp. Administrativas": [-x for x in dre_data["Despesas Administrativas"]],
                    "(-) Despesas RH": [-x for x in dre_data["Despesas RH"]],
                    "(=) EBITDA": [
                        dre_data["Receita"][i] 
                        - dre_data["Impostos Sobre Venda"][i]
                        - dre_data["Despesas Operacionais"][i] 
                        - dre_data["Despesas Administrativas"][i]
                        - dre_data["Despesas RH"][i]
                        for i in range(len(anos))
                    ],
                    "(-) Desp. Extra Op.": [-x for x in dre_data["Despesas Extra Operacional"]],
                    "(-) Dividendos": [-x for x in dre_data["Dividendos"]],
                    "(-) Imp. s/ Resultado": [-x for x in dre_data["Impostos Sobre Resultado"]],
                    "(=) FLUXO LÍQUIDO": dre_data["Lucro Líquido"]
                }
                
                fluxo_df = pd.DataFrame(fluxo_data)
                fluxo_df.index = [f"Ano {ano}" for ano in anos]
                criar_slide_com_tabela(f"{emoji} Fluxo de Caixa - {cenario}", fluxo_df)

        # ========== RECOMENDAÇÕES ESTRATÉGICAS FINAIS ==========
        
        # Análise de risco consolidada
        risco_text = """ANÁLISE DE RISCOS E OPORTUNIDADES

🔍 PRINCIPAIS RISCOS IDENTIFICADOS:

📉 RISCOS FINANCEIROS:
• Margem líquida abaixo de 15% em alguns cenários
• Dependência de financiamentos externos
• Volatilidade dos preços de commodities
• Custos operacionais crescentes

🌡️ RISCOS OPERACIONAIS:
• Variabilidade climática e safras
• Pragas e doenças nas culturas
• Disponibilidade de mão de obra
• Logística e armazenamento

💹 RISCOS DE MERCADO:
• Flutuação cambial
• Demanda internacional
• Concorrência regional
• Políticas governamentais

🎯 OPORTUNIDADES ESTRATÉGICAS:

✅ MELHORIAS OPERACIONAIS:
• Tecnologia agrícola (precisão, automação)
• Diversificação de culturas
• Integração vertical
• Sustentabilidade e certificações

📈 CRESCIMENTO:
• Expansão de área cultivada
• Novos mercados e canais
• Produtos com maior valor agregado
• Parcerias estratégicas

⚡ AÇÕES RECOMENDADAS:
• Monitoramento contínuo de indicadores
• Planos de contingência para cenários adversos
• Investimento em tecnologia e capacitação
• Diversificação como estratégia de risco"""
        
        criar_slide_texto("⚠️ Gestão de Riscos e Oportunidades", risco_text)

        # Slide de recomendações finais
        recomendacoes_text = f"""RECOMENDAÇÕES ESTRATÉGICAS FINAIS

🎯 PRIORIDADES IMEDIATAS (0-6 MESES):

✅ GESTÃO FINANCEIRA:
• Implementar controle rigoroso de custos
• Renegociar condições de financiamentos
• Estabelecer reservas de emergência
• Monitorar fluxo de caixa semanalmente

🌱 OPERAÇÕES AGRÍCOLAS:
• Otimizar uso de insumos por cultura
• Implementar técnicas de agricultura de precisão
• Diversificar mix de culturas conforme análise
• Investir em capacitação da equipe

📊 MONITORAMENTO:
• Acompanhar indicadores mensalmente
• Comparar com benchmarks do setor
• Ajustar cenários conforme realidade
• Revisar estratégias trimestralmente

🔄 MÉDIO PRAZO (6-18 MESES):

📈 CRESCIMENTO SUSTENTÁVEL:
• Expandir culturas com melhor performance
• Investir em tecnologia e equipamentos
• Desenvolver novos canais de comercialização
• Buscar certificações e selos de qualidade

⚖️ LONGO PRAZO (18+ MESES):

🏭 INTEGRAÇÃO E INOVAÇÃO:
• Considerar integração vertical
• Desenvolver produtos com valor agregado  
• Explorar mercados internacionais
• Implementar práticas ESG

💡 PRÓXIMOS PASSOS:
1. Definir KPIs de acompanhamento
2. Estabelecer comitê de análise mensal
3. Criar dashboard de indicadores
4. Planejar revisões semestrais do plano"""
        
        criar_slide_texto("🚀 Plano de Ação Estratégico", recomendacoes_text)

        # SLIDE FINAL: CONCLUSÕES
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "📊 RELATÓRIO COMPLETO FINALIZADO"
        subtitle.text = f"""✅ Análise Financeira Abrangente Concluída

📋 CONTEÚDO GERADO:
• {len(prs.slides)} slides com análises detalhadas
• Indicadores de {len(nomes_cenarios)} cenários completos
• Análises individuais por cultura
• Pareceres técnicos especializados  
• Fluxos de caixa projetados
• Recomendações estratégicas

🔧 BASE TÉCNICA:
• Sistema 5_Indicadores.py
• Dados periodo {anos[0]}-{anos[-1]}
• Metodologia financeira consolidada
• Pareceres baseados em benchmarks do agronegócio

📅 Gerado em: {datetime.now().strftime('%d/%m/%Y às %H:%M')}
🏢 Sistema de Gestão de Plantio
💼 Análise Profissional Completa"""

        # Salvar apresentação
        output_ppt = BytesIO()
        prs.save(output_ppt)
        output_ppt.seek(0)
        
        st.write(f"✅ PowerPoint COMPLETO gerado com {len(prs.slides)} slides!")
        st.write("📋 Incluindo TODAS as análises do 5_Indicadores.py:")
        st.write("   • DREs consolidados e por cultura")
        st.write("   • Indicadores financeiros detalhados")
        st.write("   • Pareceres técnicos especializados")
        st.write("   • Fluxos de caixa projetados")
        st.write("   • Comparativos entre cenários")
        st.write("   • Análises individuais por cultura")
        st.write("   • Recomendações estratégicas")
        
        return output_ppt
        
    except Exception as e:
        st.error(f"❌ Erro ao gerar PowerPoint: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        return None


def criar_relatorio_ppt(all_indicators, all_dre_data, df_culturas_for_excel, nomes_cenarios, anos, all_indicators_cultura_cenarios=None):
    """
    Função compatível que chama a versão melhorada
    """
    return criar_relatorio_ppt_completo(all_indicators, all_dre_data, df_culturas_for_excel, nomes_cenarios, anos, all_indicators_cultura_cenarios)
